from __future__ import annotations

import csv
import json
import zipfile
from io import StringIO
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from src.embedded_extract import extract_all_embeddings
from src.ooxml_media import copy_media_bundle
from src.options import ConvertOptions
from src.post_assets import append_post_sections
from src.text_formatting import (
    shape_is_visible,
    slide_title_and_body,
    strip_known_footer_lines,
    text_frame_to_markdown,
)


def _ensure_dirs_artifact(out_dir: Path) -> tuple[Path, Path]:
    assets = out_dir / "assets"
    for sub in ("images", "charts", "tables", "media", "other"):
        (assets / sub).mkdir(parents=True, exist_ok=True)
    doc = out_dir / "document.md"
    return doc, assets


def _table_to_gfm(table) -> str:
    rows: list[list[str]] = []
    for r in range(len(table.rows)):
        row = table.rows[r]
        cells = []
        for c in range(len(row.cells)):
            t = (row.cells[c].text or "").replace("|", "\\|").replace("\n", " ")
            cells.append(t)
        rows.append(cells)
    if not rows:
        return ""
    width = max(len(r) for r in rows)
    norm = [r + [""] * (width - len(r)) for r in rows]
    header = norm[0]
    sep = ["---"] * width
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(sep) + " |",
    ]
    for row in norm[1:]:
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


def _table_to_csv_file(table, path: Path) -> None:
    rows: list[list[str]] = []
    for r in range(len(table.rows)):
        row = table.rows[r]
        rows.append([(row.cells[c].text or "") for c in range(len(row.cells))])
    with path.open("w", newline="", encoding="utf-8") as f:
        w = csv.writer(f)
        w.writerows(rows)


def _chart_title(chart) -> str | None:
    try:
        if chart.has_title and chart.chart_title.has_text_frame:
            t = chart.chart_title.text_frame.text
            return (t or "").strip() or None
    except Exception:
        pass
    return None


def _chart_series_csv(chart) -> str | None:
    """Best-effort CSV from first plot's series (labels + values)."""
    try:
        plot = chart.plots[0]
        lines: list[list[str]] = []
        for series in plot.series:
            name = str(series.name) if series.name is not None else "Series"
            vals = []
            try:
                vals = [str(v) for v in series.values]
            except Exception:
                vals = []
            lines.append([name] + vals)
        if not lines:
            return None
        buf = StringIO()
        w = csv.writer(buf)
        for row in lines:
            w.writerow(row)
        return buf.getvalue()
    except Exception:
        return None


def _write_minimal_png(path: Path) -> None:
    from PIL import Image

    Image.new("RGBA", (4, 4), (240, 240, 240, 255)).save(path, format="PNG")


def convert_pptx(
    input_path: Path,
    output_path: Path,
    options: ConvertOptions,
) -> None:
    input_path = input_path.resolve()
    prs = Presentation(str(input_path))

    if options.artifact_layout:
        out_dir = output_path.resolve()
        out_dir.mkdir(parents=True, exist_ok=True)
        doc_path, assets_dir = _ensure_dirs_artifact(out_dir)
        images_dir = assets_dir / "images"
        manifest: list[dict] = []
        media_orphans: list[str] = []
        with zipfile.ZipFile(input_path, "r") as zf:
            _, media_orphans = copy_media_bundle(
                zf,
                assets_dir / "media",
                verbose=options.verbose,
            )
        extract_all_embeddings(input_path, assets_dir, options, manifest, verbose=options.verbose)
    else:
        doc_path = output_path.resolve()
        doc_path.parent.mkdir(parents=True, exist_ok=True)
        images_dir = (options.images_dir or doc_path.parent / "images").resolve()
        images_dir.mkdir(parents=True, exist_ok=True)
        assets_dir = None
        manifest = None
        media_orphans = []

    lines: list[str] = []
    img_counter: dict[int, int] = {}
    table_counter: dict[int, int] = {}
    chart_counter: dict[int, int] = {}

    for si, slide in enumerate(prs.slides, start=1):
        title, shapes = slide_title_and_body(slide)
        if options.artifact_layout:
            lines.append(f"## Slide {si}: {title}\n")
        else:
            if si == 1 and options.title_slide_h1:
                lines.append(f"# {title}\n")
            else:
                lines.append(f"## {title}\n")

        for shape in shapes:
            if not shape_is_visible(shape):
                continue
            ph = img_counter.get(si, 0)
            th = table_counter.get(si, 0)
            ch = chart_counter.get(si, 0)

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                ph += 1
                img_counter[si] = ph
                ext = shape.image.ext or "png"
                fname = f"slide_{si}_img_{ph}.{ext}"
                fpath = images_dir / fname
                fpath.write_bytes(shape.image.blob)
                rel = fpath.resolve().relative_to(doc_path.parent.resolve())
                lines.append(f"![Slide {si} image {ph}]({rel.as_posix()})\n")

            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                th += 1
                table_counter[si] = th
                t = shape.table
                lines.append(_table_to_gfm(t))
                lines.append("")
                if options.artifact_layout and options.table_csv:
                    csv_path = assets_dir / "tables" / f"slide_{si}_table_{th}.csv"
                    _table_to_csv_file(t, csv_path)

            elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                ch += 1
                chart_counter[si] = ch
                chart = shape.chart
                ct = _chart_title(chart)
                if ct:
                    lines.append(f"- Chart: {ct}\n")
                else:
                    lines.append("<!-- Chart: title not available -->\n")
                if options.artifact_layout and assets_dir is not None:
                    if options.chart_data:
                        csv_body = _chart_series_csv(chart)
                        if csv_body:
                            cp = assets_dir / "charts" / f"slide_{si}_chart_{ch}.csv"
                            cp.write_text(csv_body, encoding="utf-8")
                    if options.chart_image:
                        pngp = assets_dir / "charts" / f"slide_{si}_chart_{ch}.png"
                        _write_minimal_png(pngp)

            elif shape.has_text_frame:
                md = text_frame_to_markdown(shape.text_frame)
                if md:
                    lines.append(md)
                    lines.append("")

        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            if notes_tf is not None:
                note_text = (notes_tf.text or "").strip()
                if note_text:
                    for ln in note_text.splitlines():
                        lines.append(f"> {ln}")
                    lines.append("")

    body = "\n".join(lines).strip() + "\n"
    if options.strip_known_footers:
        body = strip_known_footer_lines(body) + "\n"

    doc_path.write_text(body, encoding="utf-8")

    if options.artifact_layout and assets_dir is not None:
        if media_orphans:
            extra = ["## Extracted package media (unreferenced)\n"]
            for o in media_orphans:
                extra.append(f"- `assets/media/{o}`")
            doc_path.write_text(
                doc_path.read_text(encoding="utf-8").rstrip() + "\n\n" + "\n".join(extra) + "\n",
                encoding="utf-8",
            )
        if manifest is not None:
            man_path = assets_dir / "extraction_manifest.json"
            man_path.write_text(json.dumps({"embeddings": manifest}, indent=2), encoding="utf-8")

        append_post_sections(doc_path, assets_dir, options)
