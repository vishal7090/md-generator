"""Build a minimal .pptx for regression tests (python tests/build_fixture.py)."""

from __future__ import annotations

from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


def main() -> int:
    root = Path(__file__).resolve().parent
    out = root / "fixtures" / "minimal.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1 — title-only layout (clean H1 in classic mode)
    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    s1.shapes.title.text = "Fixture Title"
    if len(s1.placeholders) > 1:
        s1.placeholders[1].text_frame.text = ""

    # Slide 2 — blank: bullets, table, picture, notes
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    body = s2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(2)).text_frame
    body.text = "Bullet one"
    body.add_paragraph().text = "Bullet two"
    body.add_paragraph().level = 1
    body.paragraphs[-1].text = "Nested bullet"

    rows, cols = 2, 2
    left, top, width, height = Inches(0.5), Inches(3), Inches(5), Inches(1.5)
    table = s2.shapes.add_table(rows, cols, left, top, width, height).table
    table.cell(0, 0).text = "H1"
    table.cell(0, 1).text = "H2"
    table.cell(1, 0).text = "A"
    table.cell(1, 1).text = "B"

    png = _tiny_png()
    pic_stream = BytesIO(png)
    s2.shapes.add_picture(pic_stream, Inches(6), Inches(1), width=Inches(2))

    notes = s2.notes_slide.notes_text_frame
    notes.text = "Speaker note line"

    prs.save(str(out))
    print(f"Wrote {out}")
    return 0


def _tiny_png() -> bytes:
    from PIL import Image

    buf = BytesIO()
    Image.new("RGB", (8, 8), color=(200, 100, 50)).save(buf, format="PNG")
    return buf.getvalue()


if __name__ == "__main__":
    raise SystemExit(main())
