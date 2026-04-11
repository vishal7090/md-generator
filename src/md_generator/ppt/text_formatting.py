from __future__ import annotations

import re

from pptx.oxml.ns import qn
from pptx.shapes.base import BaseShape
from pptx.slide import Slide


def _paragraph_numbered(p_el) -> bool:
    ppr = p_el.find(qn("a:pPr"))
    if ppr is None:
        return False
    return ppr.find(qn("a:buAutoNum")) is not None or ppr.find(qn("a:buChar")) is not None


def text_frame_to_markdown(text_frame) -> str:
    """Convert a python-pptx text frame to Markdown (bullets / simple numbering)."""
    lines: list[str] = []
    numbered_counters: dict[int, int] = {}

    for para in text_frame.paragraphs:
        raw = (para.text or "").replace("\r", "").strip()
        if not raw:
            continue
        level = int(para.level or 0)
        indent = "  " * level
        p_el = para._p
        numbered = _paragraph_numbered(p_el)
        if numbered:
            n = numbered_counters.get(level, 0) + 1
            numbered_counters[level] = n
            for k in list(numbered_counters):
                if k > level:
                    numbered_counters[k] = 0
            lines.append(f"{indent}{n}. {raw}")
        else:
            for k in list(numbered_counters):
                if k >= level:
                    numbered_counters[k] = 0
            lines.append(f"{indent}- {raw}")

    return "\n".join(lines)


KNOWN_FOOTER_PATTERNS = (
    re.compile(r"^\s*confidential\s*$", re.I),
    re.compile(r"^\s*proprietary\s*$", re.I),
    re.compile(r"^\s*internal use only\s*$", re.I),
)


def strip_known_footer_lines(text: str) -> str:
    out: list[str] = []
    for line in text.splitlines():
        stripped = line.strip()
        if any(p.match(stripped) for p in KNOWN_FOOTER_PATTERNS):
            continue
        out.append(line)
    return "\n".join(out).strip()


def _is_slide_title_shape(slide: Slide, shape: BaseShape, title_shape: BaseShape | None) -> bool:
    if title_shape is None:
        return False
    if shape is title_shape:
        return True
    try:
        return getattr(shape, "shape_id", None) == getattr(title_shape, "shape_id", None)
    except Exception:
        return False


def slide_title_and_body(slide: Slide) -> tuple[str, list[BaseShape]]:
    """Return (title, other_shapes_in_order) for content extraction."""
    title = ""
    title_shape = slide.shapes.title
    if title_shape is not None and title_shape.has_text_frame:
        title = (title_shape.text or "").strip()
    if not title:
        title = "Untitled"

    ordered: list[BaseShape] = []
    for shape in slide.shapes:
        if _is_slide_title_shape(slide, shape, title_shape):
            continue
        ordered.append(shape)
    return title, ordered


def shape_is_visible(shape: BaseShape) -> bool:
    try:
        if getattr(shape, "visible", True) is False:
            return False
    except Exception:
        pass
    return True
